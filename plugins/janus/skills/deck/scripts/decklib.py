#!/usr/bin/env python3
"""decklib — build a PowerPoint deck *natively* on top of an existing template.

The idea: keep a template's master / theme / layouts / logo / footer, throw away
its sample slides, and rebuild your own slides on the template's own layouts so
the result looks like it was made in that template.

Typical use (see SKILL.md for the full flow):

    from decklib import Deck, RGB
    d = Deck("template.pptx")          # ea_font defaults to "Noto Sans CJK JP"
    d.master_replace_text("v0.0-TODO", "My Lab")   # fix template chrome text
    d.strip_slides(keep=lambda s: "Thank you" in d.slide_text(s))  # keep 1 slide

    s = d.add("TITLE")                 # add a slide on the layout named TITLE
    d.text(s, 0, "My Title", bold=True, size=32)
    d.fit(s, 0, 2.3, 2.2, 9.4, 1.3)    # always set FULL geometry when moving a ph

    s = d.add("CUSTOM_4_17_1")
    d.body(s, 2, [("Head", "detail line"), ("Head2", "detail2")])
    d.prose(s, 3, "Paragraph one.\\n\\nParagraph two.")   # narrative text, no bullets
    d.table(s, 1.0, 2.7, 5.4, 3.0, ("Col A", "Col B"),
            [("a", "1"), ("b", "2")], title="My table")
    d.picture(s, "diagram.png", 1.0, 2.3, width=11.4)
    d.svg(s, "diagram.svg", 1.0, 2.3, width=5.0, light=True)  # SVG in one call
    d.refs(s, [("OpenShift Docs", "https://docs.redhat.com/x")])  # call LAST on the slide

    d.move_to_end(lambda s: "Thank you" in d.slide_text(s))
    d.save("out.pptx")

Discovery helpers (run first to learn a template's vocabulary):
    Deck("template.pptx").describe_layouts()
    Deck("template.pptx").describe_slides()
"""
import os
import tempfile
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE


def RGB(hexstr):
    """RGB('EE0000') -> RGBColor."""
    hexstr = hexstr.lstrip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


class Deck:
    def __init__(self, template, ea_font="Noto Sans CJK JP"):
        # ea_font: East-Asian typeface forced on every run we create. On Linux
        # "Noto Sans CJK JP" avoids LibreOffice falling back to Chinese glyphs.
        self.prs = Presentation(template)
        self.ea_font = ea_font

    # ---------------- discovery ----------------
    def _layouts(self):
        for m in self.prs.slide_masters:
            for l in m.slide_layouts:
                yield l

    def layout(self, name):
        for l in self._layouts():
            if l.name == name:
                return l
        raise KeyError(f"no layout named {name!r}; run describe_layouts()")

    def describe_layouts(self):
        seen = set()
        for l in self._layouts():
            if l.name in seen:
                continue
            seen.add(l.name)
            print(f"\n### LAYOUT {l.name!r}")
            inch = lambda v: (v or 0) / 914400.0  # EMU -> inches
            for p in l.placeholders:
                pf = p.placeholder_format
                print(f"   idx={pf.idx} type={pf.type} "
                      f"pos=({inch(p.left):.1f},{inch(p.top):.1f}) "
                      f"size=({inch(p.width):.1f}x{inch(p.height):.1f})")

    def slide_text(self, slide):
        return " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)

    def describe_slides(self):
        for i, s in enumerate(self.prs.slides):
            hidden = s._element.get("show") == "0"
            txt = self.slide_text(s).strip().replace("\n", " / ")[:70]
            print(f"[{i:3}] {'(hidden)' if hidden else 'VISIBLE ':8} "
                  f"layout={s.slide_layout.name:16} :: {txt}")

    # ---------------- slide-set surgery ----------------
    def strip_slides(self, keep=lambda slide: False, keep_first=False):
        """Delete the template's sample slides; keep those where keep(slide) is True.

        keep_first=True keeps only the FIRST matching slide (templates often hide
        several duplicate sample slides — e.g. multiple "Thank you" pages)."""
        lst = self.prs.slides._sldIdLst
        kept = False
        for slide, sldId in list(zip(self.prs.slides, list(lst))):
            if keep(slide) and not (keep_first and kept):
                kept = True
                continue
            self.prs.part.drop_rel(sldId.get(qn("r:id")))
            lst.remove(sldId)

    def move_to_end(self, predicate):
        lst = self.prs.slides._sldIdLst
        for slide, sldId in list(zip(self.prs.slides, list(lst))):
            if predicate(slide):
                lst.remove(sldId)
                lst.append(sldId)
                return

    def add(self, layout_name):
        return self.prs.slides.add_slide(self.layout(layout_name))

    def master_replace_text(self, find, replace):
        """Rewrite a fixed text on the master (e.g. a version/footer placeholder)."""
        for m in self.prs.slide_masters:
            for sh in m.shapes:
                if sh.has_text_frame and find in sh.text_frame.text:
                    for p in sh.text_frame.paragraphs:
                        for r in p.runs:
                            if find in r.text:
                                r.text = r.text.replace(find, replace)

    # ---------------- run/text helpers ----------------
    def ph(self, slide, idx):
        for p in slide.placeholders:
            if p.placeholder_format.idx == idx:
                return p
        return None

    def _set_ea(self, run):
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", self.ea_font)

    def _style(self, run, size=None, bold=None, color=None):
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        self._set_ea(run)

    def text(self, slide, idx, s, size=None, bold=None, color=None):
        """Set a single-line placeholder's text."""
        p = self.ph(slide, idx)
        if p is None:
            return None
        tf = p.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = s
        self._style(r, size=size, bold=bold, color=color)
        return p

    def body(self, slide, idx, items, head_sz=14, desc_sz=12, gap=8,
             head_color=None, desc_color=None):
        """Fill a body placeholder with (head, detail) pairs as level0/level1 bullets.

        head is bold (level 0); detail (optional) is an indented level-1 line.
        """
        head_color = head_color or RGB("151515")
        desc_color = desc_color or RGB("595959")
        p = self.ph(slide, idx)
        tf = p.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for head, desc in items:
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            para.level = 0
            if not first:
                para.space_before = Pt(gap)
            r = para.add_run(); r.text = head
            self._style(r, size=head_sz, bold=True, color=head_color)
            if desc:
                d = tf.add_paragraph(); d.level = 1
                r = d.add_run(); r.text = desc
                self._style(r, size=desc_sz, bold=False, color=desc_color)
            first = False
        return p

    def _no_bullet(self, para):
        pPr = para._p.get_or_add_pPr()
        if pPr.find(qn("a:buNone")) is not None:
            return
        bu = pPr.makeelement(qn("a:buNone"), {})
        ref = pPr.find(qn("a:defRPr"))
        if ref is not None:
            ref.addprevious(bu)
        else:
            pPr.append(bu)

    def prose(self, slide, idx, text, size=12, bold=None, color=None, gap=8):
        """Fill a placeholder with free paragraphs — no bullet glyphs.

        Blank lines ('\\n\\n') split spaced paragraphs; a single '\\n' stays a
        line break inside the paragraph. Use instead of body() for narrative
        text (quotes, official statements): body() renders a blank line as an
        empty bullet."""
        color = color or RGB("151515")
        p = self.ph(slide, idx)
        tf = p.text_frame
        tf.clear()
        tf.word_wrap = True
        blocks = [b for b in (s.strip("\n") for s in text.split("\n\n")) if b.strip()]
        for i, block in enumerate(blocks):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if i:
                para.space_before = Pt(gap)
            for j, line in enumerate(block.split("\n")):
                if j:
                    para.add_line_break()
                r = para.add_run(); r.text = line
                self._style(r, size=size, bold=bold, color=color)
            self._no_bullet(para)
        return p

    def disclaimer(self, slide, idx, conditions, notes=(), size=12,
                   note_size=10, gap=6, color=None, note_color=None):
        """Fill a body placeholder as a disclaimer slide: conditions as
        bullets, notes as smaller grey non-bulleted ※-lines below."""
        color = color or RGB("151515")
        note_color = note_color or RGB("8C8C8C")
        p = self.ph(slide, idx)
        tf = p.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for c in conditions:
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                para.space_before = Pt(gap)
            r = para.add_run(); r.text = c
            self._style(r, size=size, color=color)
            first = False
        for i, n in enumerate(notes):
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                para.space_before = Pt(gap * 2 if i == 0 else gap)
            r = para.add_run(); r.text = n if n.startswith("※") else "※ " + n
            self._style(r, size=note_size, color=note_color)
            self._no_bullet(para)
            first = False
        return p

    def fit(self, slide, idx, left, top, width, height):
        """Reposition a placeholder. ALWAYS pass full geometry — setting only
        width/top on an inherited placeholder leaves the others at 0 and the box
        flies off-slide or collapses to zero width."""
        p = self.ph(slide, idx)
        p.left, p.top = Inches(left), Inches(top)
        p.width, p.height = Inches(width), Inches(height)
        return p

    move = fit  # alias

    def clear(self, slide, idx):
        p = self.ph(slide, idx)
        if p is not None:
            p.text_frame.clear()
        return p

    # ---------------- pictures & tables ----------------
    def picture(self, slide, src, left, top, width=None, height=None):
        """src may be a path or raw image bytes."""
        if isinstance(src, (bytes, bytearray)):
            src = BytesIO(src)
        kw = {}
        if width is not None:
            kw["width"] = Inches(width)
        if height is not None:
            kw["height"] = Inches(height)
        return slide.shapes.add_picture(src, Inches(left), Inches(top), **kw)

    def svg(self, slide, src, left, top, width=None, height=None,
            px_width=2000, light=False, palette=None, background="white"):
        """Insert an SVG (file path or '<svg…>' markup) as a picture in one
        call — renders to PNG via svgtools.render_svg (rsvg-convert/inkscape)
        instead of the manual SVG → rsvg-convert → picture() dance.
        light=True remaps the SVG's hex colours to the light template palette
        first (svgtools.auto_light_map; palette overrides its entries).
        Needs svgtools.py importable next to this file."""
        import svgtools
        text = src if src.lstrip().startswith("<") else open(src, encoding="utf-8").read()
        if light:
            text = svgtools.recolor_svg(
                text, svgtools.auto_light_map(svgtools.svg_colors(text), palette))
        fd, tmp = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        try:
            svgtools.render_svg(text, tmp, width=px_width, background=background)
            with open(tmp, "rb") as f:
                data = f.read()
        finally:
            os.unlink(tmp)
        return self.picture(slide, data, left, top, width=width, height=height)

    # ---------------- reference footnotes ----------------
    def _text_bottom(self, sh):
        """Estimate the bottom (inches) of the *rendered* text in a
        text-frame shape. Placeholder boxes often stretch to the slide
        bottom while their text stops much higher — using the box bottom
        would leave no room for footnotes."""
        tf = sh.text_frame
        if not tf.text.strip():
            return 0.0
        width_in = (sh.width or 0) / 914400.0 or 1.0
        h = 0.0
        for para in tf.paragraphs:
            sizes = [r.font.size.pt for r in para.runs if r.font.size] or [14]
            size = max(sizes)
            text = "".join(r.text for r in para.runs)
            cpl = max(int(width_in * 72 / (size * 0.7)), 8)
            h += max(1, -(-len(text) // cpl)) * size * 1.5 / 72.0
            if para.space_before:
                h += para.space_before.pt / 72.0
        return (sh.top or 0) / 914400.0 + h

    def content_bottom(self, slide):
        """Lowest bottom edge (inches) over the slide's content — everything
        below it is the safe zone for footnotes. Text shapes count with
        their estimated rendered-text bottom (not the full placeholder box);
        pictures/tables with their full box."""
        bottom = 0.0
        for sh in slide.shapes:
            if sh.has_text_frame:
                bottom = max(bottom, self._text_bottom(sh))
            elif sh.top is not None and sh.height is not None:
                bottom = max(bottom, (sh.top + sh.height) / 914400.0)
        return bottom

    def refs(self, slide, items, size=9, color=None, left=0.5, width=None,
             bottom_margin=0.28, compact_at=3, gap=0.08):
        """Reference footnote at the slide bottom. Call LAST on the slide:
        it measures content_bottom() and fits the refs below the content.
        With compact_at or more items — or when one-line-per-ref doesn't fit
        above bottom_margin — the refs are joined into a single wrapped
        paragraph one point smaller, instead of overlapping the body (the
        2-column full-height-body overlap problem).

        items: strings, or (label, url) pairs joined as "label — url"."""
        color = color or RGB("8C8C8C")
        strs = [it if isinstance(it, str) else " — ".join(str(x) for x in it if x)
                for it in items]
        if not strs:
            return None
        emu = 914400.0
        slide_w = self.prs.slide_width / emu
        slide_h = self.prs.slide_height / emu
        if width is None:
            # stop short of the bottom-right footer chrome (version/logo)
            width = max(slide_w - left - 2.9, slide_w * 0.6)
        floor = self.content_bottom(slide) + gap
        line_h = size * 1.45 / 72.0
        if len(strs) >= compact_at or floor + len(strs) * line_h > slide_h - bottom_margin:
            size = max(size - 1, 7)
            line_h = size * 1.45 / 72.0
            joined = "   |   ".join(strs)
            # over-estimate wrap (0.7 em/char, mixed CJK+ASCII) so the box
            # grows upward rather than spilling off-slide
            chars_per_line = max(int(width * 72 / (size * 0.7)), 20)
            est_h = -(-len(joined) // chars_per_line) * line_h
            texts = [joined]
        else:
            est_h = len(strs) * line_h
            texts = strs
        # prefer just above the bottom margin; move up if content intrudes,
        # but never start so low the refs run off-slide
        top = min(max(floor, slide_h - bottom_margin - est_h),
                  slide_h - est_h - 0.02)
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(est_h + 0.05))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        for i, t in enumerate(texts):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = para.add_run(); r.text = t
            self._style(r, size=size, bold=False, color=color)
        return tb

    def table(self, slide, x, y, w, h, header, rows, title=None,
              header_bg="EE0000", header_fg="FFFFFF", zebra="F5F5F5",
              body_fg="151515", first_col_ratio=0.62, size=11):
        if title:
            cap = slide.shapes.add_textbox(Inches(x), Inches(y - 0.35), Inches(w), Inches(0.3))
            r = cap.text_frame.paragraphs[0].add_run(); r.text = title
            self._style(r, size=12, bold=True, color=RGB(header_bg))
        gf = slide.shapes.add_table(len(rows) + 1, len(header), Inches(x), Inches(y),
                                    Inches(w), Inches(h))
        tbl = gf.table
        if len(header) == 2:
            tbl.columns[0].width = int(Inches(w) * first_col_ratio)
            tbl.columns[1].width = int(Inches(w) * (1 - first_col_ratio))

        def cell(c, txt, bold, color, bg):
            c.margin_top = Pt(3); c.margin_bottom = Pt(3)
            c.margin_left = Pt(7); c.margin_right = Pt(5)
            if bg is None:
                c.fill.background()
            else:
                c.fill.solid(); c.fill.fore_color.rgb = RGB(bg)
            c.text_frame.word_wrap = True
            r = c.text_frame.paragraphs[0].add_run(); r.text = txt
            self._style(r, size=size, bold=bold, color=color)

        for ci, htext in enumerate(header):
            cell(tbl.cell(0, ci), htext, True, RGB(header_fg), header_bg)
        for ri, row in enumerate(rows, start=1):
            bg = zebra if ri % 2 == 0 else "FFFFFF"
            for ci, val in enumerate(row):
                cell(tbl.cell(ri, ci), val, False, RGB(body_fg), bg)
        return gf

    def add_textbox(self, slide, left, top, width, height, text,
                    size=None, bold=None, color=None, align=None):
        """Add a free-floating text box (not a placeholder) at the given inches.

        Use when no spare placeholder is available — e.g. a date on a cover whose
        TITLE layout has no free placeholder. `align` takes a PP_ALIGN value
        (e.g. PP_ALIGN.CENTER); defaults to the layout/left default."""
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        if align is not None:
            para.alignment = align
        r = para.add_run()
        r.text = text
        self._style(r, size=size, bold=bold, color=color)
        return tb

    # ---------------- output ----------------
    def save(self, path):
        self.prs.save(path)
        return path
